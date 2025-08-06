from pptx import Presentation
from pptx.util import Inches

# Initialize presentation
presentation = Presentation()

# Helper function to add a slide
def add_slide(presentation, title, content, placeholder=None):
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout
    slide.shapes.title.text = title
    content_box = slide.placeholders[1]
    content_box.text = content
    if placeholder:
        content_box.text += f"\n\n[Placeholder: {placeholder}]"
    return slide

# Slide 1: Title Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
slide.shapes.title.text = "Automatic Number Plate Recognition Using Deep Learning"
slide.placeholders[1].text = "Leveraging Python, YOLOv8, and EasyOCR for Detection and Reading"

# Slide 2: Introduction
add_slide(
    presentation,
    "Introduction",
    "The project focuses on detecting vehicles, extracting license plates, and reading numbers. "
    "Real-world applications include traffic management, toll automation, and law enforcement.",
    placeholder="Infographic of applications"
)

# Slide 3: Objectives
add_slide(
    presentation,
    "Objectives",
    "- Detect and locate number plates in images/videos using YOLOv8.\n"
    "- Recognize and extract text using EasyOCR.\n"
    "- Ensure high accuracy and real-time processing.",
    placeholder="Conceptual diagram of detection and OCR integration"
)

# Slide 4: Tools and Libraries
add_slide(
    presentation,
    "Tools and Libraries",
    "Technologies Used:\n- Python for programming\n- YOLOv8 for object detection\n"
    "- EasyOCR for text recognition\n- OpenCV for image processing\n\n"
    "Advantages: YOLOv8’s speed and EasyOCR’s multilingual capabilities.",
    placeholder="Logos of tools"
)

# Slide 5: Dataset Overview
add_slide(
    presentation,
    "Dataset Overview",
    "The dataset includes bounding box annotations for cars and license plates. "
    "It has two versions: `test.csv` for raw results and `test_interpolated.csv` for smoother tracking.",
    placeholder="Table of dataset statistics"
)

# Slide 6: System Architecture
add_slide(
    presentation,
    "System Architecture",
    "Workflow:\n1. Input image or video\n2. YOLOv8 for vehicle and plate detection\n"
    "3. Extract plate regions\n4. OCR for reading license numbers.",
    placeholder="Workflow diagram"
)

# Slide 7: YOLOv8 for Object Detection
add_slide(
    presentation,
    "YOLOv8 for Object Detection",
    "YOLOv8 is used for its speed and accuracy. Training data and performance metrics are highlighted. "
    "Sample images include bounding boxes for detected cars and plates.",
    placeholder="Performance metrics graph"
)

# Slide 8: EasyOCR for Text Extraction
add_slide(
    presentation,
    "EasyOCR for Text Extraction",
    "EasyOCR reads text from detected plate regions. Challenges include poor image quality and obstructed plates. "
    "Examples show successful and failed OCR readings.",
    placeholder="Images with OCR outputs"
)

# Slide 9: Code Implementation
add_slide(
    presentation,
    "Code Implementation",
    "The project is implemented across multiple scripts:\n- Detection (`main.py`)\n- Data Processing (`add_missing_data.py`)\n"
    "- Visualization (`visualize.py`)\n- Tracking (`sort.py`).\n\nKey snippets are highlighted.",
    placeholder="Annotated code snippets"
)

# Slide 10: Visualization of Results
add_slide(
    presentation,
    "Visualization of Results",
    "Results include annotated images/videos with bounding boxes and license numbers. Confidence scores and OCR outputs are highlighted.",
    placeholder="Comparison images and outputs table"
)

# Slide 11: Evaluation and Performance
add_slide(
    presentation,
    "Evaluation and Performance",
    "Metrics include detection accuracy, OCR accuracy, and system performance. Challenges include false positives and dataset limitations.",
    placeholder="Bar chart of performance metrics"
)

# Slide 12: Applications
add_slide(
    presentation,
    "Applications",
    "Use cases:\n- Traffic enforcement\n- Toll automation systems\n- Parking management\n\nFuture Enhancements:\n"
    "- Multilingual plates\n- Better low-light detection.",
    placeholder="Icons with captions for use cases"
)

# Slide 13: Conclusion
add_slide(
    presentation,
    "Conclusion",
    "Summary:\n- Successfully integrated YOLOv8 and EasyOCR for real-time recognition.\n"
    "- Challenges and potential improvements identified.\n\nTakeaway: Importance of this technology in various industries.",
    placeholder="Key achievements as bullet points"
)

# Slide 14: Q&A / Thank You
add_slide(
    presentation,
    "Questions?",
    "Thank you for your attention!\nContact: [Your contact details]",
    placeholder="Relevant image or workflow diagram"
)

# Save the presentation
presentation.save("ANPR_Presentation.pptx")
